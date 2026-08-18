import os
from pathlib import Path


def parse_file(file_path: str) -> str:
    """解析文件，返回纯文本内容"""
    ext = Path(file_path).suffix.lower()

    if ext == ".txt" or ext == ".md":
        return parse_txt(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def parse_txt(file_path: str) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse_pdf(file_path: str) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n".join(texts)


def parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    return "\n".join(texts)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """
    按段落切分，尽量保持语义完整。
    段落过短则合并，过长则按句号截断。
    """
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    if not paragraphs:
        return []

    chunks = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) <= chunk_size:
            current = (current + "\n" + para).strip() if current else para
        else:
            if current:
                chunks.append(current)
            # 如果单段落超过chunk_size，按句号截断
            if len(para) > chunk_size:
                sub_chunks = split_long_paragraph(para, chunk_size, overlap)
                chunks.extend(sub_chunks)
                current = ""
            else:
                current = para
    if current:
        chunks.append(current)

    # 添加overlap
    if overlap > 0 and len(chunks) > 1:
        overlapped = []
        for i, chunk in enumerate(chunks):
            if i > 0:
                prev_tail = chunks[i - 1][-overlap:] if len(chunks[i - 1]) > overlap else chunks[i - 1]
                chunk = prev_tail + "\n" + chunk
            overlapped.append(chunk)
        return overlapped

    return chunks


def parse_xlsx(file_path: str) -> str:
    """解析 xlsx 文件为纯文本（用于模糊匹配资料，非题库导入）"""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        texts.append(f"【{sheet_name}】")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
            if row_text:
                texts.append(row_text)
    wb.close()
    return "\n".join(texts)


def split_long_paragraph(text: str, chunk_size: int, overlap: int) -> list[str]:
    sentences = []
    for s in text.replace("。", "。|").replace("！", "！|").replace("？", "？|").split("|"):
        s = s.strip()
        if s:
            sentences.append(s)

    chunks = []
    current = ""
    for sent in sentences:
        if len(current) + len(sent) <= chunk_size:
            current = (current + sent).strip() if current else sent
        else:
            if current:
                chunks.append(current)
            if len(sent) > chunk_size:
                # Very long sentence, force split
                for i in range(0, len(sent), chunk_size - overlap):
                    chunks.append(sent[i:i + chunk_size])
                current = ""
            else:
                current = sent
    if current:
        chunks.append(current)
    return chunks
