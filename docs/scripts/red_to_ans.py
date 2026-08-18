#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
red_to_ans.py — 多格式题库预处理脚本

把教材类资料中"红色字体"的文字转成 【ANS:xxx】 标记，
用于送给 DP / ChatGPT 等大模型识别"哪些是填空答案"。

支持格式（按文件后缀自动识别）:
    .docx           Word 2007+  (依赖 python-docx)
    .doc            Word 97-2003 (Windows + Word 自动转 .docx 后处理)
    .pdf            PDF 文档    (依赖 PyMuPDF)
    .xlsx           Excel       (依赖 openpyxl)
    .html .htm      HTML 网页   (依赖 beautifulsoup4)
    .pptx           PowerPoint  (依赖 python-pptx)
    .rtf            RTF 富文本   (依赖 striprtf + 简易颜色解析)

使用:
    python red_to_ans.py <input_file> <output.txt>
    python red_to_ans.py 教材.docx 教材_marked.txt
    python red_to_ans.py 资料.pdf  资料_marked.txt
    python red_to_ans.py 课件.pptx 课件_marked.txt

安装依赖（按需）:
    pip install python-docx PyMuPDF openpyxl beautifulsoup4 python-pptx striprtf

输出说明:
    所有红色文字会被包裹成 【ANS:原文】 形式
    其它内容保持原文（保留段落分隔）
    输出 UTF-8 .txt 文件
"""

import os
import sys
import shutil
import tempfile
from pathlib import Path


# ============================================
# 红色判定（统一颜色识别规则）
# ============================================

def is_red_rgb(r, g, b):
    """判断 RGB 值是否为红色系。

    判定规则：
      - R ≥ 150（红通道足够强）
      - G ≤ 100, B ≤ 100（绿蓝通道不能过亮）
      - R 必须明显大于 G 和 B（差值 ≥ 50）
    覆盖：纯红 (255,0,0) / 深红 (200,30,30) / 玫红 (200,50,80) 等
    不会误判：黑色、灰色、棕色、橙黄
    """
    try:
        r, g, b = int(r), int(g), int(b)
    except (TypeError, ValueError):
        return False
    return r >= 150 and g <= 100 and b <= 100 and r > g + 50 and r > b + 50


def is_red_hex(hex_str):
    """十六进制颜色（#RRGGBB / RRGGBB / 0xRRGGBB）→ 是否红色"""
    if not hex_str:
        return False
    s = str(hex_str).strip().lstrip('#').lstrip('0x').lstrip('0X')
    if len(s) == 3:
        s = ''.join(c * 2 for c in s)
    if len(s) != 6:
        return False
    try:
        return is_red_rgb(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return False


def is_red_named(name):
    """常见英文/中文颜色名 → 是否红色"""
    if not name:
        return False
    s = str(name).strip().lower()
    return s in {
        'red', 'darkred', 'firebrick', 'crimson', 'maroon', 'indianred',
        '红色', '深红', '红', '朱红', '紫红',
    }


# ============================================
# 处理器：各格式独立函数，统一输出 list[str] 段落
# ============================================

def process_docx(path):
    """处理 .docx（Word 2007+）"""
    try:
        from docx import Document
        from docx.shared import RGBColor
    except ImportError:
        sys.exit("缺少依赖: pip install python-docx")

    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        buf = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            red = False
            try:
                if run.font.color and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    if isinstance(rgb, RGBColor):
                        red = is_red_rgb(rgb[0], rgb[1], rgb[2])
            except Exception:
                pass
            if red and text.strip():
                buf.append(f'【ANS:{text.strip()}】')
            else:
                buf.append(text)
        line = ''.join(buf).strip()
        if line:
            lines.append(line)

    # 同时处理表格内的红字（教材常用表格列答案）
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    cbuf = []
                    for run in para.runs:
                        text = run.text
                        if not text:
                            continue
                        red = False
                        try:
                            if run.font.color and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                red = is_red_rgb(rgb[0], rgb[1], rgb[2])
                        except Exception:
                            pass
                        cbuf.append(f'【ANS:{text.strip()}】' if red and text.strip() else text)
                    cline = ''.join(cbuf).strip()
                    if cline:
                        lines.append(cline)
    return lines


def process_doc(path):
    """处理 .doc（Word 97-2003）— Windows 上用 win32com 自动转 .docx 再处理"""
    if sys.platform != 'win32':
        sys.exit(
            "✗ .doc 格式需要 Windows + Microsoft Word\n"
            "  → 请手动用 Word/WPS 另存为 .docx 后再处理\n"
            "  → 或在 macOS/Linux 用 libreoffice --convert-to docx 转换"
        )
    try:
        import win32com.client as w
    except ImportError:
        sys.exit("缺少依赖: pip install pywin32")

    tmp_docx = tempfile.mktemp(suffix='.docx')
    word = None
    try:
        word = w.Dispatch('Word.Application')
        word.Visible = False
        abs_src = os.path.abspath(path)
        d = word.Documents.Open(abs_src)
        d.SaveAs(tmp_docx, FileFormat=16)  # 16 = docx
        d.Close()
    finally:
        if word:
            try: word.Quit()
            except Exception: pass

    try:
        return process_docx(tmp_docx)
    finally:
        try: os.remove(tmp_docx)
        except Exception: pass


def process_pdf(path):
    """处理 .pdf — 用 PyMuPDF 提取文字 + 颜色信息"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        sys.exit("缺少依赖: pip install PyMuPDF")

    lines = []
    doc = fitz.open(path)
    for page in doc:
        # 用 dict 模式获取每个 span 的颜色
        blocks = page.get_text("dict").get("blocks", [])
        for block in blocks:
            if block.get("type", 0) != 0:  # 只处理文本块
                continue
            for line in block.get("lines", []):
                buf = []
                for span in line.get("spans", []):
                    text = span.get("text", "")
                    if not text:
                        continue
                    # color 是整数，例如 0xFF0000
                    color_int = span.get("color", 0)
                    r = (color_int >> 16) & 0xFF
                    g = (color_int >> 8) & 0xFF
                    b = color_int & 0xFF
                    red = is_red_rgb(r, g, b)
                    if red and text.strip():
                        buf.append(f'【ANS:{text.strip()}】')
                    else:
                        buf.append(text)
                ln = ''.join(buf).strip()
                if ln:
                    lines.append(ln)
    doc.close()
    return lines


def process_xlsx(path):
    """处理 .xlsx — 用 openpyxl 检查每个单元格的字体颜色

    特点：Excel 单元格颜色信息有两类
      1. 单元格全体字色（一致颜色）
      2. 单元格内富文本不同段不同色（罕见）
    我们主要处理第 1 种，覆盖 90%+ 教材表格场景。
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        sys.exit("缺少依赖: pip install openpyxl")

    wb = load_workbook(path, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows():
            row_buf = []
            for cell in row:
                v = cell.value
                if v is None or str(v).strip() == '':
                    continue
                text = str(v)
                red = False
                try:
                    color = cell.font.color
                    if color and color.type == 'rgb' and color.rgb:
                        # color.rgb 形如 'FFFF0000'，前 2 位是 alpha
                        hex_str = color.rgb[-6:] if len(color.rgb) >= 6 else color.rgb
                        red = is_red_hex(hex_str)
                except Exception:
                    pass
                if red:
                    row_buf.append(f'【ANS:{text.strip()}】')
                else:
                    row_buf.append(text)
            if row_buf:
                lines.append(' | '.join(row_buf))
    return lines


def process_html(path):
    """处理 .html / .htm — BeautifulSoup 解析"""
    try:
        from bs4 import BeautifulSoup, NavigableString
    except ImportError:
        sys.exit("缺少依赖: pip install beautifulsoup4")

    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        html = f.read()
    soup = BeautifulSoup(html, 'html.parser')

    # 删除脚本/样式
    for tag in soup(['script', 'style', 'meta', 'link']):
        tag.decompose()

    def is_red_element(el):
        """检查元素是否被设置为红色（font color / style color）"""
        # <font color="red"> 或 <font color="#FF0000">
        if el.name == 'font' and el.get('color'):
            c = el['color']
            return is_red_named(c) or is_red_hex(c)
        # style="color:red" / style="color:#FF0000"
        style = el.get('style', '')
        if style:
            import re
            m = re.search(r'color\s*:\s*([^;]+)', style, re.IGNORECASE)
            if m:
                c = m.group(1).strip()
                if c.startswith('rgb'):
                    nums = re.findall(r'\d+', c)
                    if len(nums) >= 3:
                        return is_red_rgb(nums[0], nums[1], nums[2])
                return is_red_named(c) or is_red_hex(c)
        return False

    lines = []

    def walk(node, in_red=False):
        if isinstance(node, NavigableString):
            text = str(node).strip()
            if text:
                if in_red:
                    lines.append(f'【ANS:{text}】')
                else:
                    lines.append(text)
            return
        if not hasattr(node, 'name') or node.name is None:
            return
        # 段落级元素后加换行符
        block_tags = {'p', 'div', 'br', 'tr', 'li', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6'}
        red_here = in_red or is_red_element(node)
        for child in node.children:
            walk(child, red_here)
        if node.name in block_tags:
            lines.append('\n')

    walk(soup)
    # 合并连续换行
    raw = ''.join(lines)
    cleaned = []
    for line in raw.split('\n'):
        line = line.strip()
        if line:
            cleaned.append(line)
    return cleaned


def process_pptx(path):
    """处理 .pptx — python-pptx 提取每个 run 的颜色"""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
    except ImportError:
        sys.exit("缺少依赖: pip install python-pptx")

    prs = Presentation(path)
    lines = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_buf = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                buf = []
                for run in para.runs:
                    text = run.text
                    if not text:
                        continue
                    red = False
                    try:
                        if run.font.color and run.font.color.type is not None:
                            rgb = run.font.color.rgb
                            if rgb:
                                red = is_red_rgb(rgb[0], rgb[1], rgb[2])
                    except Exception:
                        pass
                    if red and text.strip():
                        buf.append(f'【ANS:{text.strip()}】')
                    else:
                        buf.append(text)
                ln = ''.join(buf).strip()
                if ln:
                    slide_buf.append(ln)
        if slide_buf:
            lines.append(f'--- 幻灯片 {slide_idx} ---')
            lines.extend(slide_buf)
    return lines


def process_rtf(path):
    """处理 .rtf — 简易颜色解析（不依赖额外库）

    RTF 内部用 \cf{n} 切换前景色，颜色表在文件头 {\colortbl ;\red255\green0\blue0; ...}
    实现一个最小可用版本，覆盖常见教材 RTF。
    """
    import re
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        rtf = f.read()

    # 解析 colortbl
    m = re.search(r'\{\\colortbl\s*([^}]*)\}', rtf)
    color_table = [None]  # 索引 0 默认是黑/无色
    if m:
        for cdef in re.finditer(r'\\red(\d+)\\green(\d+)\\blue(\d+)', m.group(1)):
            color_table.append((int(cdef.group(1)), int(cdef.group(2)), int(cdef.group(3))))

    # 简易 token 解析，按 \cf 切色，按 \par 切段
    lines = []
    buf = []
    in_red = False
    pos = 0
    tokens = re.split(r'(\\[A-Za-z]+\d*\s?|[\{\}])', rtf[rtf.find('}') + 1:])
    for tok in tokens:
        if not tok:
            continue
        if tok.startswith('\\cf'):
            try:
                idx = int(tok[3:].strip())
                if idx < len(color_table) and color_table[idx]:
                    r, g, b = color_table[idx]
                    in_red = is_red_rgb(r, g, b)
                else:
                    in_red = False
            except ValueError:
                in_red = False
            continue
        if tok in ('\\par', '\\line'):
            line = ''.join(buf).strip()
            if line:
                lines.append(line)
            buf = []
            continue
        if tok.startswith('\\') or tok in ('{', '}'):
            continue  # 忽略其它控制字
        if tok.strip():
            if in_red:
                buf.append(f'【ANS:{tok.strip()}】')
            else:
                buf.append(tok)
    if buf:
        line = ''.join(buf).strip()
        if line:
            lines.append(line)
    return lines


# ============================================
# 主入口
# ============================================

HANDLERS = {
    '.docx': process_docx,
    '.doc':  process_doc,
    '.pdf':  process_pdf,
    '.xlsx': process_xlsx,
    '.html': process_html,
    '.htm':  process_html,
    '.pptx': process_pptx,
    '.rtf':  process_rtf,
}


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    src = sys.argv[1]
    dst = sys.argv[2]

    if not os.path.exists(src):
        sys.exit(f"✗ 文件不存在: {src}")

    ext = Path(src).suffix.lower()
    handler = HANDLERS.get(ext)
    if handler is None:
        sys.exit(
            f"✗ 不支持的文件格式: {ext}\n"
            f"  支持: {', '.join(HANDLERS.keys())}"
        )

    print(f'正在处理 [{ext}]: {src}')
    try:
        lines = handler(src)
    except SystemExit:
        raise
    except Exception as e:
        sys.exit(f"✗ 处理失败: {e}")

    # 统计 ANS 标记数
    ans_count = sum(line.count('【ANS:') for line in lines)
    total_chars = sum(len(l) for l in lines)

    text = '\n'.join(lines)
    with open(dst, 'w', encoding='utf-8') as f:
        f.write(text)

    print(f'✓ 处理完成')
    print(f'  段落数 : {len(lines)}')
    print(f'  字符数 : {total_chars}')
    print(f'  红字标记: {ans_count} 处')
    print(f'  输出   : {dst}')

    if ans_count == 0:
        print()
        print('⚠ 没有检测到任何红色字体')
        print('  可能原因：')
        print('    1. 文档里压根没用红字 → 这个文件不需要预处理，直接发 DP')
        print('    2. 用了"高亮/底色"而非"字体颜色"')
        print('    3. 颜色不是标准红（如紫红/深紫）')
        print('    用 Word 检查：选中疑似答案文字 → 看"开始 → 字体颜色"是不是红')


if __name__ == '__main__':
    main()
